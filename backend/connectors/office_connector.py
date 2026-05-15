"""Office connector stubs for Excel, Word, and PowerPoint.
Attempts to use openpyxl, python-docx and python-pptx if available; falls back to simple file outputs otherwise.
"""
import os

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import docx
except Exception:
    docx = None

try:
    import pptx
except Exception:
    pptx = None


def create_excel(path: str, rows: list) -> bool:
    """Create a simple Excel file from rows (list of lists). Returns True on success."""
    try:
        if openpyxl:
            wb = openpyxl.Workbook()
            ws = wb.active
            for r in rows:
                ws.append(r)
            wb.save(path)
            return True
        else:
            # fallback: write CSV
            import csv
            with open(path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerows(rows)
            return True
    except Exception:
        return False


def create_word(path: str, paragraphs: list) -> bool:
    """Create a simple Word doc from paragraphs. Returns True on success."""
    try:
        if docx:
            d = docx.Document()
            for p in paragraphs:
                d.add_paragraph(str(p))
            d.save(path)
            return True
        else:
            with open(path, 'w', encoding='utf-8') as f:
                for p in paragraphs:
                    f.write(str(p) + "\n\n")
            return True
    except Exception:
        return False


def create_ppt(path: str, slides: list) -> bool:
    """Create a simple PPTX with slides list of titles. Returns True on success."""
    try:
        if pptx:
            prs = pptx.Presentation()
            for title in slides:
                slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = str(title)
            prs.save(path)
            return True
        else:
            with open(path, 'w', encoding='utf-8') as f:
                for s in slides:
                    f.write(str(s) + "\n---\n")
            return True
    except Exception:
        return False


if __name__ == '__main__':
    # quick manual smoke test
    create_excel('sample.xlsx', [["A","B"],[1,2]])
    create_word('sample.docx', ["Hello", "World"]) 
    create_ppt('sample.pptx', ["Slide 1", "Slide 2"])